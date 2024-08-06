import uuid
from flask import Flask, request, jsonify
from flask_cors import CORS
from werkzeug.utils import secure_filename
import os
import json
from pypdf import PdfReader
from pptx import Presentation
import google.generativeai as genai
from dotenv import load_dotenv

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*", "allow_headers": "*"}})

load_dotenv()

api_key = os.getenv("API_KEY")
genai.configure(api_key=api_key)
model_gen = genai.GenerativeModel('gemini-1.5-flash')

UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = [page.extract_text() for page in reader.pages]
    return " ".join(text)

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = [
        shape.text for slide in presentation.slides
        for shape in slide.shapes if hasattr(shape, "text")
    ]
    return "\n".join(text)

def extract_text(file_path):
    if file_path.lower().endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.lower().endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file format")

def generate_flashcards(text, number_of_cards):
    response = model_gen.generate_content(
        f"Generate exactly {number_of_cards} flashcards from the following text. Each flashcard should be a JSON object with two properties: 'front' and 'back'. The 'front' property should contain the question, and the 'back' property should contain the answer. The result must be a single JSON array with all flashcard objects inside it. Each object should be separated by a comma, and there should be no trailing comma after the last object. Ensure that the JSON array is properly formatted without any extra text, code blocks, or delimiters. Use double quotes for all keys and string values. Here's the text: {text}"
    )

    result = response.text if hasattr(response, 'text') else response
    
    print("Raw result:", result)
    
    cleaned_data = result.strip()

    # Clean up code block delimiters if they exist
    if cleaned_data.startswith('```json'):
        cleaned_data = cleaned_data.replace('```json', '').replace('```', '').strip()
    elif cleaned_data.startswith('```javascript') or cleaned_data.startswith('```js'):
        cleaned_data = cleaned_data.replace('```javascript', '').replace('```js', '').replace('```', '').strip()

    print("Cleaned data:", cleaned_data)

    try:
        flashcards = json.loads(cleaned_data)
        print("Parsed flashcards:", flashcards)
    except json.JSONDecodeError as e:
        print(f"Error parsing JSON: {e}")
        flashcards = []

    return flashcards



@app.route('/create-flashcards', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if file:
        # Generate a UUID for the file
        unique_id = str(uuid.uuid4())
        original_filename = secure_filename(file.filename)
        extension = os.path.splitext(original_filename)[1]
        unique_filename = f"{unique_id}{extension}"
        
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        file.save(file_path)

        try:
            number_of_cards = int(request.form.get('number_of_cards', 10))
            text = extract_text(file_path)
            flashcards = generate_flashcards(text, number_of_cards)
            os.remove(file_path)
            return jsonify({'flashcards': flashcards}), 200
        except Exception as e:
            return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(port=5000)
